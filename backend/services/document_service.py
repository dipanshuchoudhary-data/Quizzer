import pdfplumber
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from pathlib import Path

from backend.integrations.ocr.tesseract import extract_text_from_image


# ----------------------------------------
# PDF Extraction (Text-based)
# ----------------------------------------

def extract_pdf_text(path: str) -> str:
    text_parts = []

    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

    return "\n".join(text_parts).strip()


# ----------------------------------------
# Scanned PDF (Safe OCR Fallback)
# ----------------------------------------

def extract_scanned_pdf_text(path: str) -> str:
    text_parts = []
    doc = fitz.open(path)

    try:
        for page_index in range(len(doc)):
            page = doc.load_page(page_index)
            pix = page.get_pixmap()

            image_path = f"{path}_page_{page_index}.png"
            pix.save(image_path)

            try:
                ocr_text = extract_text_from_image(image_path)
                if ocr_text:
                    text_parts.append(ocr_text)
            finally:
                Path(image_path).unlink(missing_ok=True)

    finally:
        doc.close()

    return "\n".join(text_parts).strip()


# ----------------------------------------
# DOCX Extraction
# ----------------------------------------

def extract_docx_text(path: str) -> str:
    doc = DocxDocument(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text]
    return "\n".join(paragraphs).strip()


# ----------------------------------------
# PPTX Extraction
# ----------------------------------------

def extract_pptx_text(path: str) -> str:
    prs = Presentation(path)
    text_parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)

    return "\n".join(text_parts).strip()


# ----------------------------------------
# Unified Dispatcher
# ----------------------------------------

def extract_text_from_file(path: str, file_type: str) -> str:

    file_type = file_type.lower()

    if file_type == "pdf":
        # Try native PDF text extraction first
        text = extract_pdf_text(path)

        # Fallback to OCR only if empty
        if not text:
            text = extract_scanned_pdf_text(path)

        return text

    if file_type == "docx":
        return extract_docx_text(path)

    if file_type == "pptx":
        return extract_pptx_text(path)

    if file_type in ["png", "jpg", "jpeg"]:
        return extract_text_from_image(path)

    raise ValueError(f"Unsupported file type: {file_type}")
